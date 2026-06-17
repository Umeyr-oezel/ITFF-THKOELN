"""Assemble the professional board deck from the bare charts (Aufgabe 1).

Builds a clean 16:9 PowerPoint: a title slide plus four content slides with
native titles, a consistent brand background, the transparent TH Koeln logo,
a footer and slide numbers. The charts are the *_bare.png variants (no baked
title/logo), so every slide carries exactly one title and one logo.

Run after generating the bare charts:
    PRES_BARE=1 python analysis/insider_insights.py
    PRES_BARE=1 python analysis/architecture_seaborn.py
    PRES_BARE=1 python analysis/architecture_diagrams.py
    python analysis/build_presentation.py
"""
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import config  # noqa: E402

PRES_DIR = os.path.join(config.OUTPUT_DIR, "presentation")
OUT_FILE = os.path.join(PRES_DIR, "Presentation - ITFF - redesigned.pptx")
FONT = "Arial"

# Accurate per-chart source: SEC data charts vs. our own diagrams/codebase.
SOURCES = {
    "sentiment_all_time": "Data: SEC EDGAR, Form 4 (2020-2025) - own analysis",
    "trend_all_time": "Data: SEC EDGAR, Form 4 (2020-2025) - own analysis",
    "covid_buyers": "Data: SEC EDGAR, Form 4 (March 2020) - own analysis",
    "system_pipeline": "Data: SEC EDGAR, Form 4 (2020-2025) - own analysis",
    "system_architecture": "Own representation",
    "code_architecture": "Own representation",
    "code_layers": "Own analysis - project codebase (lines of code)",
}

BG = RGBColor(0xFA, 0xFA, 0xFA)
TEXT = RGBColor(0x2D, 0x2D, 0x2D)
SUB = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x1A, 0x6B, 0x54)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU = 914400  # per inch


def chart(name):
    """Path to a bare chart by stem (e.g. 'sentiment_all_time')."""
    return os.path.join(PRES_DIR, f"{name}_bare.png")


def text(slide, s, l, t, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box with our brand font settings."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = s
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = FONT
    run.font.color.rgb = color
    return tb


def accent(slide, l, t, w, h, color=GREEN):
    """A thin filled rectangle used as a design accent / rule."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                                 Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def background(slide):
    """Fill the whole slide with the brand background colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def logo(slide, top=0.32, height=0.62):
    """Drop the transparent logo into the top-right corner."""
    if not os.path.isfile(config.LOGO_PATH):
        return
    iw, ih = Image.open(config.LOGO_PATH).size
    w = height * iw / ih
    slide.shapes.add_picture(config.LOGO_PATH, Inches(13.333 - 0.45 - w),
                             Inches(top), height=Inches(height))


def src_of(path):
    """Look up the correct source caption for a chart by its file stem."""
    stem = os.path.basename(path).replace("_bare.png", "")
    return SOURCES.get(stem, "Own analysis")


def place(slide, img, box):
    """Fit an image inside (l, t, w, h) keeping aspect; add a source caption."""
    l, t, bw, bh = box
    iw, ih = Image.open(img).size
    ar = iw / ih
    if bw / bh > ar:
        h = bh
        w = h * ar
    else:
        w = bw
        h = w / ar
    il, it = l + (bw - w) / 2, t + (bh - h) / 2
    slide.shapes.add_picture(img, Inches(il), Inches(it), width=Inches(w))
    text(slide, src_of(img), il, it + h + 0.04, w, 0.3, 8.5, SUB,
         align=PP_ALIGN.CENTER)


def footer(slide, page):
    """Source note (left) and slide number (right)."""
    text(slide, "Group 01 - IT for Finance  -  TH Koeln",
         0.45, 7.05, 8, 0.35, 9, SUB)
    text(slide, f"{page} / 4", 11.5, 7.05, 1.35, 0.35, 9, SUB,
         align=PP_ALIGN.RIGHT)


def content_slide(prs, page, title, subtitle, images):
    """One content slide: accent + title + subtitle + chart(s) + footer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    accent(slide, 0.45, 0.45, 0.9, 0.07)
    text(slide, title, 0.45, 0.55, 11, 0.7, 28, TEXT, bold=True)
    text(slide, subtitle, 0.47, 1.18, 11.5, 0.45, 14, SUB)
    logo(slide)

    area = (0.45, 1.72, 12.43, 4.8)  # l, t, w, h (room left for captions)
    if len(images) == 1:
        place(slide, images[0], area)
    else:
        gap = 0.3
        cw = (area[2] - gap) / 2
        place(slide, images[0], (area[0], area[1], cw, area[3]))
        place(slide, images[1], (area[0] + cw + gap, area[1], cw, area[3]))
    footer(slide, page)
    return slide


def title_slide(prs):
    """Opening slide: project title, hook, team, logo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    accent(slide, 0.85, 2.35, 1.6, 0.09)
    text(slide, "Insider Trading on SEC Form 4", 0.8, 2.5, 11.7, 1.0, 40, TEXT,
         bold=True)
    text(slide, "What 4.6 million insider filings reveal about conviction  -  2020-2025",
         0.83, 3.55, 11.5, 0.6, 18, SUB)
    text(slide, "Group 01  -  IT for Finance", 0.83, 5.4, 11, 0.4, 14, GREEN, bold=True)
    text(slide, "Emil  -  Umeyr  -  Kenan  -  Matthias        |        TH Koeln, June 2026",
         0.83, 5.85, 11.5, 0.4, 12, SUB)
    logo(slide, top=0.45, height=0.75)
    return slide


def main():
    """Build and save the redesigned deck."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    content_slide(
        prs, 1, "When the market crashed, insiders bought",
        "March 2020: the insider buy-share hit a six-year high (#1 of 72 months)",
        [chart("sentiment_all_time")],
    )
    content_slide(
        prs, 2, "System Architecture",
        "A six-phase, idempotent data pipeline - 4.6M rows from SEC EDGAR",
        [chart("system_architecture"), chart("system_pipeline")],
    )
    content_slide(
        prs, 3, "Code Architecture",
        "Clean layers, a swappable database, 30 unit tests",
        [chart("code_architecture"), chart("code_layers")],
    )
    content_slide(
        prs, 4, "What the data revealed",
        "Selling is the norm (3.3x) - except when conviction strikes",
        [chart("trend_all_time"), chart("covid_buyers")],
    )

    prs.save(OUT_FILE)
    print(f"saved {OUT_FILE}")


if __name__ == "__main__":
    main()
